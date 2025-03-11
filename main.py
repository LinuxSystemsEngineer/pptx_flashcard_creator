import os
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def create_directories():
    os.makedirs("./output", exist_ok=True)

def clear_screen():
    os.system('cls' if os.name == 'nt' else 'clear')

def print_menu():
    clear_screen()
    print("\n" + "=" * 50)
    print("        🎓 PowerPoint Flashcard Creator 🎓")
    print("=" * 50)
    print("  1️⃣  Create a new set of flashcards")
    print("  2️⃣  Exit")
    print("=" * 50)

def print_header(text):
    print("\n" + "-" * 50)
    print(f"  {text}")
    print("-" * 50)

def set_font(text_frame, text, size):
    text_frame.clear()
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

def create_flashcards():
    create_directories()
    prs = Presentation()
    output_path = "./output/output.pptx"
    set_counter = 1
    
    while True:
        print_header(f"Creating a New Flashcard Set {set_counter:03}")
        question = input("📖 Enter the question: ")
        slide_layout = prs.slide_layouts[5]  # Title Only layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        
        set_font(title.text_frame, "Question:", 34)
        p = title.text_frame.add_paragraph()
        run = p.add_run()
        run.text = "\n" + question
        run.font.name = "Arial"
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        answer = input("✅ Enter the answer: ")
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        
        set_font(title.text_frame, "Answer:", 34)
        p = title.text_frame.add_paragraph()
        run = p.add_run()
        run.text = "\n" + answer
        run.font.name = "Arial"
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        print("\n✔ Flashcard added successfully!")
        set_counter += 1
        choice = input("➕ Would you like to create another set? (1️⃣ Yes | 2️⃣ No): ")
        if choice == "2":
            break
    
    prs.save(output_path)
    print(f"📂 Flashcards saved to {output_path}")

def main():
    create_directories()
    while True:
        print_menu()
        choice = input("👉 Select an option: ")
        
        if choice == "1":
            create_flashcards()
        elif choice == "2":
            print("👋 Exiting program. Have a great day!")
            break
        else:
            print("❌ Invalid choice, please try again.")

if __name__ == "__main__":
    main()
